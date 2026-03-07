"""Generate Bee v4.0 Capabilities — PPT, DOCX, PDF"""

import os
from pathlib import Path

OUT = Path(r"d:\Projects\Workitulife\docs\presentations")
OUT.mkdir(parents=True, exist_ok=True)

# ── Shared content ──────────────────────────────────────────────

TITLE = "Bee v4.0 — AI Chief of Staff"
SUBTITLE = "Capabilities Overview"
AUTHOR = "Workitu Tech"
DATE = "March 2026"

SLIDES = [
    {
        "title": "What is Bee?",
        "bullets": [
            "AI Chief of Staff for Yonatan Perlin, Founder of Workitu Tech",
            "Operates across 5 surfaces: Claude Code, Claude.ai, Telegram, Vault API, Autonomous Agent",
            "27+ agent tools, 23 MCP operations, 19 AI skills",
            "Shared brain repo synced via GitHub",
            "Mission: reach ₪7,000/mo revenue → $1B by 2087",
        ],
    },
    {
        "title": "Architecture",
        "bullets": [
            "Central brain repo (workitu-bee-brain) on GitHub",
            "Claude Code — full filesystem, git, terminal, MCP access (primary operator)",
            "Claude.ai — memory, dashboard, skills, MCP access",
            "Telegram Bot — always-on chat with AI + all 27 tools",
            "Vault REST API — HTTPS endpoints for note management",
            "Autonomous Agent — daily cron on VPS, DeepSeek AI",
        ],
    },
    {
        "title": "Surface 1: Autonomous Agent",
        "bullets": [
            "Runs daily at 08:00 IST on Hetzner VPS",
            "DeepSeek AI (deepseek-chat / deepseek-reasoner)",
            "9 tool modules with 27 tools total",
            "ClickUp: overdue tasks, due today, comments",
            "Gmail & Calendar: unread emails, today's events",
            "Site Monitor: health checks, broken links, alerts",
            "Content Manager: blog publishing, portfolio, revalidation",
            "Performance Reporter: weekly stats + email reports",
            "Obsidian Vault: search, read, create, update notes",
        ],
    },
    {
        "title": "Surface 2: Telegram Bot",
        "bullets": [
            "Always-on via PM2 on VPS",
            "7 slash commands: /tasks, /calendar, /status, /briefing, etc.",
            "Free-form AI chat with DeepSeek + full tool access",
            "Conversation context: 20 messages, 30-min TTL",
            "Rate limited: 10 messages/minute",
            "Admin-only authentication via Telegram user ID",
        ],
    },
    {
        "title": "Surface 3: Vault REST API",
        "bullets": [
            "URL: https://api.workitu.com/api/vault/",
            "FastAPI on port 3003, behind nginx with HTTPS",
            "8 endpoints: health, search, read, create, update, list, daily, link",
            "API key authentication (X-API-Key header)",
            "Rate limit: 30 requests/minute",
            "Swagger UI + OpenAPI schema included",
        ],
    },
    {
        "title": "Surface 4: Claude Code (MCP)",
        "bullets": [
            "Primary operator — full filesystem, git, terminal access",
            "ClickUp MCP: 10+ operations (search, create, update tasks & docs)",
            "Gmail MCP: search, read, draft, send emails",
            "Google Calendar MCP: events, scheduling, free time",
            "Native: SSH to VPS, script execution, project deployment",
            "Deploys to Vercel and Hetzner VPS",
        ],
    },
    {
        "title": "Surface 5: Claude.ai Skills",
        "bullets": [
            "19 skills across 7 categories",
            "Core: Chief of Staff, Memory Engine, Skill Developer",
            "Business: Profit Mode, Outreach Composer, Client Manager",
            "Development: Fullstack Builder, PRD-to-Deploy, Site Enhancer",
            "Personal: Medical Tracker, Weekly Review",
            "Automation: Daily Briefing, YouTube Creator",
            "Storage: Obsidian Vault with [[wikilinks]]",
        ],
    },
    {
        "title": "Infrastructure",
        "bullets": [
            "VPS: Hetzner CX23 (65.109.230.136), Ubuntu 24.04",
            "Services: nginx, Telegram bot, Vault API, bee-backend, n8n",
            "Website: Vercel (blaze-post.com / workitu.com)",
            "DNS: Cloudflare",
            "External APIs: DeepSeek, ClickUp, Gmail, Calendar, Resend, Telegram",
            "Local: Lenovo IdeaPad Flex 5, Ryzen 7, 16GB, Ollama offline",
        ],
    },
    {
        "title": "Security",
        "bullets": [
            "Path validation — no directory traversal",
            "Shell injection prevention — subprocess list-form",
            "API key auth on Vault REST API",
            "Telegram admin-only via user ID check",
            "Rate limiting: 30 req/min API, 10 msg/min Telegram",
            "HTTPS with Let's Encrypt auto-renewal",
            ".env files on VPS only — never committed to repo",
        ],
    },
    {
        "title": "By the Numbers",
        "bullets": [
            "5 operational surfaces",
            "27 agent tools (autonomous)",
            "23 MCP operations (Claude Code + Claude.ai)",
            "19 AI skills (Claude.ai)",
            "8 REST API endpoints",
            "7 Telegram slash commands",
            "9 tool modules on the autonomous agent",
            "Daily automated briefings at 8:00 AM IST",
        ],
    },
]


# ── 1. PowerPoint ───────────────────────────────────────────────

def make_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color palette
    BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
    ACCENT = RGBColor(0xFF, 0xBE, 0x0B)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)

    def set_slide_bg(slide, color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_text_box(slide, left, top, width, height, text, font_size=18,
                     color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                         Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = alignment
        return tf

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, 1, 1.5, 11, 1.5, TITLE,
                 font_size=44, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 3.2, 11, 1, SUBTITLE,
                 font_size=28, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 4.5, 11, 0.6, f"{AUTHOR}  •  {DATE}",
                 font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Content slides
    for s in SLIDES:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide, BG_DARK)

        # Title bar
        add_text_box(slide, 0.8, 0.4, 11, 0.9, s["title"],
                     font_size=32, color=ACCENT, bold=True)

        # Divider line
        from pptx.shapes.autoshape import Shape
        line = slide.shapes.add_shape(
            1,  # rectangle
            Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.04)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT
        line.line.fill.background()

        # Bullets
        txBox = slide.shapes.add_textbox(
            Inches(1.0), Inches(1.7), Inches(11), Inches(5.2)
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(s["bullets"]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"▸  {bullet}"
            p.font.size = Pt(20)
            p.font.color.rgb = WHITE
            p.space_after = Pt(12)

    # Final slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_text_box(slide, 1, 2, 11, 1.5, "Thank You",
                 font_size=48, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, 1, 3.8, 11, 1,
                 "Bee v4.0 — AI Chief of Staff\nworkitu.com  •  blaze-post.com",
                 font_size=22, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    path = OUT / "Bee_v4_Capabilities.pptx"
    prs.save(str(path))
    print(f"PPTX saved: {path}")


# ── 2. Word Document ────────────────────────────────────────────

def make_docx():
    from docx import Document
    from docx.shared import Pt, Inches, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    # Title
    title = doc.add_heading(TITLE, level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for run in title.runs:
        run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    sub = doc.add_paragraph(f"{SUBTITLE}\n{AUTHOR} — {DATE}")
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for run in sub.runs:
        run.font.size = Pt(14)
        run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    doc.add_paragraph("")  # spacer

    for section in SLIDES:
        h = doc.add_heading(section["title"], level=1)
        for run in h.runs:
            run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

        for bullet in section["bullets"]:
            p = doc.add_paragraph(bullet, style="List Bullet")
            for run in p.runs:
                run.font.size = Pt(11)

        doc.add_paragraph("")  # spacer

    # Summary table
    doc.add_heading("Quick Reference", level=1)
    table = doc.add_table(rows=9, cols=2)
    table.style = "Light Grid Accent 1"
    data = [
        ("Surfaces", "5"),
        ("Agent Tools", "27"),
        ("MCP Operations", "23"),
        ("AI Skills", "19"),
        ("REST Endpoints", "8"),
        ("Telegram Commands", "7"),
        ("Tool Modules", "9"),
        ("Autonomous Schedule", "Daily 8:00 AM IST"),
        ("Revenue Target", "₪7,000/month"),
    ]
    for i, (k, v) in enumerate(data):
        table.rows[i].cells[0].text = k
        table.rows[i].cells[1].text = v

    path = OUT / "Bee_v4_Capabilities.docx"
    doc.save(str(path))
    print(f"DOCX saved: {path}")


# ── 3. PDF ──────────────────────────────────────────────────────

def make_pdf():
    from fpdf import FPDF

    def pdf_safe(text):
        """Replace Unicode chars that latin-1 can't encode."""
        return (text
                .replace("\u2014", "-")   # em-dash
                .replace("\u2013", "-")   # en-dash
                .replace("\u2018", "'")   # left single quote
                .replace("\u2019", "'")   # right single quote
                .replace("\u201c", '"')   # left double quote
                .replace("\u201d", '"')   # right double quote
                .replace("\u2022", "*")   # bullet
                .replace("\u25b8", ">")   # triangle bullet
                .replace("\u20aa", "NIS") # shekel sign
                .replace("\u2192", "->")  # arrow
                .encode("latin-1", "replace").decode("latin-1"))

    class BeePDF(FPDF):
        def header(self):
            self.set_font("Helvetica", "B", 10)
            self.set_text_color(150, 150, 150)
            self.cell(0, 8, "Bee v4.0 - AI Chief of Staff", align="R")
            self.ln(12)

        def footer(self):
            self.set_y(-15)
            self.set_font("Helvetica", "I", 8)
            self.set_text_color(150, 150, 150)
            self.cell(0, 10, f"Workitu Tech - {DATE}  |  Page {self.page_no()}/{{nb}}", align="C")

    pdf = BeePDF()
    pdf.alias_nb_pages()
    pdf.set_auto_page_break(auto=True, margin=20)

    # Title page
    pdf.add_page()
    pdf.ln(60)
    pdf.set_font("Helvetica", "B", 36)
    pdf.set_text_color(26, 26, 46)
    pdf.cell(0, 15, pdf_safe(TITLE), align="C", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(5)
    pdf.set_font("Helvetica", "", 20)
    pdf.set_text_color(100, 100, 100)
    pdf.cell(0, 10, SUBTITLE, align="C", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(10)
    pdf.set_font("Helvetica", "", 14)
    pdf.cell(0, 8, f"{AUTHOR}  -  {DATE}", align="C", new_x="LMARGIN", new_y="NEXT")

    # Content pages
    for section in SLIDES:
        pdf.add_page()
        pdf.set_font("Helvetica", "B", 22)
        pdf.set_text_color(26, 26, 46)
        pdf.cell(0, 12, pdf_safe(section["title"]), new_x="LMARGIN", new_y="NEXT")

        # Divider
        pdf.set_draw_color(255, 190, 11)
        pdf.set_line_width(0.8)
        pdf.line(10, pdf.get_y() + 2, 200, pdf.get_y() + 2)
        pdf.ln(8)

        pdf.set_font("Helvetica", "", 12)
        pdf.set_text_color(50, 50, 50)

        for bullet in section["bullets"]:
            safe = pdf_safe(bullet)
            pdf.cell(8)
            pdf.multi_cell(0, 7, f">  {safe}", new_x="LMARGIN", new_y="NEXT")
            pdf.ln(2)

    path = OUT / "Bee_v4_Capabilities.pdf"
    pdf.output(str(path))
    print(f"PDF saved: {path}")


# ── Run ─────────────────────────────────────────────────────────

if __name__ == "__main__":
    make_pptx()
    make_docx()
    make_pdf()
    print("\nAll 3 documents generated in:", OUT)
